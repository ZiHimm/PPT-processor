from pptx import Presentation
import logging
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)

class PPTReader:
    def __init__(self, ppt_path: str):
        try:
            self.prs = Presentation(ppt_path)
        except Exception as e:
            logger.error(f"Failed to load PPT {ppt_path}: {e}")
            raise

    def get_slide_indexes(self):
        return range(len(self.prs.slides))


    def get_text_shapes(self, slide_index, debug=False):
        try:
            slide = self.prs.slides[slide_index]
        except IndexError:
            logger.warning(f"Slide {slide_index} not found")
            return []

        results = []
        
        def extract_from_shape(shape, parent_left=0, parent_top=0):
            """Recursively extract text from any shape type"""
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            
            # Handle text frames
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    results.append({
                        "text": text,
                        "left": shape.left + parent_left,
                        "top": shape.top + parent_top
                    })
            
            # Handle tables
            elif shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            results.append({
                                "text": cell_text,
                                "left": shape.left + parent_left,
                                "top": shape.top + parent_top
                            })
            
            # Handle groups (recursive)
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for sub_shape in shape.shapes:
                    extract_from_shape(sub_shape, shape.left, shape.top)
            
            # Optional: Handle other shapes with text
            elif hasattr(shape, "text") and shape.text.strip():
                results.append({
                    "text": shape.text.strip(),
                    "left": shape.left + parent_left,
                    "top": shape.top + parent_top
                })
        
        # Process all shapes on slide
        for shape in slide.shapes:
            extract_from_shape(shape)
        
        return results