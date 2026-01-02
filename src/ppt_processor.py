"""
Enhanced PPT Processor with caching, error recovery, and detailed logging
FIXED VERSION - Proper logger initialization
"""

import re
import logging
import json
import hashlib
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple, Callable
from datetime import datetime
from functools import lru_cache
from dataclasses import dataclass
import yaml
import traceback

from pptx import Presentation

# Import utils functions with fallbacks
try:
    from .utils import format_number, extract_date_from_text, clean_text, safe_int
except ImportError:
    # Define fallback functions if utils can't be imported
    def format_number(value: Any) -> Optional[float]:
        if value is None:
            return None
        if isinstance(value, (int, float)):
            return float(value)
        if isinstance(value, str):
            try:
                return float(value.replace(',', ''))
            except:
                return None
        return None
    
    def extract_date_from_text(text: str) -> Optional[str]:
        import re
        match = re.search(r'(\d{1,2}/\d{1,2}/\d{2,4}|\d{4}/\d{1,2}/\d{1,2})', text)
        return match.group(1) if match else None
    
    def clean_text(text: str) -> str:
        if not text:
            return ""
        import re
        text = re.sub(r'\s+', ' ', text)
        text = re.sub(r'[^\w\s\.\,\-\:\/]', '', text)
        return text.strip()
    
    def safe_int(value: Any, default: int = 0) -> int:
        try:
            if value is None:
                return default
            if isinstance(value, str):
                value = value.replace(',', '').replace(' ', '')
            return int(float(value))
        except:
            return default

# Import DataValidator with fallback
try:
    from .data_validator import DataValidator
except ImportError:
    # Create a dummy DataValidator if not available
    class DataValidator:
        def __init__(self, config_path: str = None):
            pass
        
        def validate_and_clean(self, data: Dict, dashboard_type: str) -> Optional[Dict]:
            return data  # Just return data as-is without validation

@dataclass
class ProcessingMetrics:
    """Track processing performance and quality"""
    total_slides: int = 0
    processed_slides: int = 0
    failed_slides: int = 0
    extracted_records: int = 0
    processing_time: float = 0.0
    data_quality_score: float = 0.0
    confidence_scores: List[float] = None
    
    def __post_init__(self):
        if self.confidence_scores is None:
            self.confidence_scores = []

class EnhancedPPTProcessor:
    """Enhanced PPT processor with caching, validation, and detailed logging"""
    
    def __init__(self, config_path: str = "config/dashboard_config.yaml"):
        # Setup logger FIRST - don't create it yet
        self._setup_logging()
        
        # Now create logger after setup
        self.logger = logging.getLogger(__name__)
        
        # Then load config
        self.config = self._load_config(config_path)
        self.patterns = self._compile_patterns()
        self.validator = DataValidator(config_path)
        self.cache_dir = Path("cache")
        self.cache_dir.mkdir(exist_ok=True)
        
        # Initialize counters
        self.metrics = ProcessingMetrics()
        self.extraction_history = []
    
    def _setup_logging(self):
        """Configure logging without basicConfig to avoid conflicts"""
        # Get or create logger
        logger = logging.getLogger(__name__)
        
        # Only add handlers if they don't already exist
        if not logger.handlers:
            log_format = '%(asctime)s - %(name)s - %(levelname)s - %(message)s'
            date_format = '%Y-%m-%d %H:%M:%S'
            
            # File handler
            try:
                file_handler = logging.FileHandler('ppt_processor.log', encoding='utf-8')
                file_handler.setLevel(logging.INFO)
                file_handler.setFormatter(logging.Formatter(log_format, date_format))
                logger.addHandler(file_handler)
            except Exception as e:
                print(f"Failed to create file handler: {e}")
            
            # Console handler
            console_handler = logging.StreamHandler()
            console_handler.setLevel(logging.INFO)
            console_handler.setFormatter(logging.Formatter(log_format, date_format))
            logger.addHandler(console_handler)
            
            logger.setLevel(logging.INFO)
    
    def _load_config(self, config_path: str) -> Dict[str, Any]:
        """Load and validate configuration"""
        try:
            with open(config_path, 'r', encoding='utf-8') as f:
                config = yaml.safe_load(f)
            
            # Validate required sections
            if config and 'dashboard_config' in config:
                required_sections = ['social_media', 'global_settings']
                dashboard_config = config.get('dashboard_config', {})
                
                for section in required_sections:
                    if section not in dashboard_config:
                        self.logger.warning(f"Missing required section: {section}")
                
                self.logger.info(f"Configuration loaded from {config_path}")
                return dashboard_config
            else:
                self.logger.warning("Empty or invalid configuration file")
                return self._get_default_config()
            
        except Exception as e:
            self.logger.error(f"Failed to load config: {e}")
            return self._get_default_config()
    
    def _get_default_config(self) -> Dict[str, Any]:
        """Return default configuration if file not found"""
        self.logger.warning("Using default configuration")
        return {
            'social_media': {'keywords': [], 'metrics': []},
            'community_marketing': {'keywords': [], 'metrics': []},
            'kol_engagement': {'keywords': [], 'metrics': []},
            'performance_marketing': {'keywords': [], 'metrics': []},
            'promotion_posts': {'keywords': [], 'metrics': []},
            'global_settings': {'date_formats': [], 'number_formats': []}
        }
    
    def _compile_patterns(self) -> Dict[str, Dict]:
        """Compile regex patterns with configuration support"""
        patterns = {
            'date': {
                'pattern': re.compile(
                    r'(?:Date|Posted|Published|Created|Updated)[:\s]*'
                    r'(\d{1,2}[/-]\d{1,2}[/-]\d{2,4}|\d{4}[/-]\d{1,2}[/-]\d{1,2}|'
                    r'\d{1,2}\s+(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\s+\d{4}|'
                    r'(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\s+\d{1,2},\s+\d{4})',
                    re.IGNORECASE
                ),
                'confidence': 0.9
            },
            'metrics': {
                'reach_views': {
                    'pattern': re.compile(r'(?:Reach|Views|Reach/Views|View Count|Reach/View)[:\s]*([\d,]+(?:\.\d+)?[KMB]?)', re.IGNORECASE),
                    'confidence': 0.85
                },
                'engagement': {
                    'pattern': re.compile(r'(?:Engagement|Engagement Rate|Total Engagement)[:\s]*([\d,]+(?:\.\d+)?[KMB]?)', re.IGNORECASE),
                    'confidence': 0.8
                },
                'likes': {
                    'pattern': re.compile(r'(?:Likes|Like Count)[:\s]*([\d,]+(?:\.\d+)?[KMB]?)', re.IGNORECASE),
                    'confidence': 0.9
                },
                'shares': {
                    'pattern': re.compile(r'(?:Shares|Share Count|Shared)[:\s]*([\d,]+(?:\.\d+)?[KMB]?)', re.IGNORECASE),
                    'confidence': 0.8
                },
                'comments': {
                    'pattern': re.compile(r'(?:Comments|Comment Count)[:\s]*([\d,]+(?:\.\d+)?[KMB]?)', re.IGNORECASE),
                    'confidence': 0.8
                },
                'saved': {
                    'pattern': re.compile(r'(?:Saved|Save Count|Saves)[:\s]*([\d,]+(?:\.\d+)?[KMB]?)', re.IGNORECASE),
                    'confidence': 0.7
                }
            },
            'entities': {
                'platform': re.compile(r'(?:Platform|Channel|Network|Page|Account)[:\s]*([A-Za-z0-9\s]+)(?:\n|$| )', re.IGNORECASE),
                'kol_name': re.compile(r'(?:KOL|Influencer|Creator|Talent)[:\s]*([A-Za-z\s]+)(?=\n|:|$| )', re.IGNORECASE),
                'community_name': re.compile(r'(?:Community|Group|Hive)[:\s]*([A-Za-z0-9\s\-]+)(?=\n|:|$)', re.IGNORECASE),
                'campaign_name': re.compile(r'(?:Campaign|Promotion|Initiative)[:\s]*([A-Za-z0-9\s\-]+)(?=\n|:|$)', re.IGNORECASE)
            }
        }
        
        return patterns
    
    def get_cache_key(self, file_path: str) -> str:
        """Generate cache key from file metadata"""
        try:
            stat = Path(file_path).stat()
            file_hash = hashlib.md5(
                f"{file_path}-{stat.st_mtime}-{stat.st_size}".encode()
            ).hexdigest()
            return file_hash[:16]
        except Exception as e:
            self.logger.warning(f"Failed to generate cache key: {e}")
            return hashlib.md5(file_path.encode()).hexdigest()[:16]
    
    def process_ppt(self, ppt_path: str, progress_callback: Callable = None) -> Dict[str, List]:
        """
        Process PPT file with caching and progress tracking
        
        Args:
            ppt_path: Path to PPT file
            progress_callback: Callback function for progress updates
            
        Returns:
            Dictionary of extracted data by dashboard type
        """
        start_time = datetime.now()
        
        # Create new metrics for this processing run
        current_metrics = ProcessingMetrics()
        
        try:
            # Check cache
            cache_key = self.get_cache_key(ppt_path)
            cache_file = self.cache_dir / f"{cache_key}.json"
            
            if cache_file.exists():
                self.logger.info(f"Loading from cache: {ppt_path}")
                with open(cache_file, 'r', encoding='utf-8') as f:
                    cached_data = json.load(f)
                    current_metrics.extracted_records = sum(len(v) for v in cached_data.values())
                    return cached_data
            
            # Process file
            self.logger.info(f"Processing: {ppt_path}")
            
            # Check if file exists
            if not Path(ppt_path).exists():
                self.logger.error(f"File not found: {ppt_path}")
                return {
                    'social_media': [],
                    'community_marketing': [],
                    'kol_engagement': [],
                    'performance_marketing': [],
                    'promotion_posts': []
                }
            
            presentation = Presentation(ppt_path)
            filename = Path(ppt_path).name
            
            results = {
                'social_media': [],
                'community_marketing': [],
                'kol_engagement': [],
                'performance_marketing': [],
                'promotion_posts': []
            }
            
            total_slides = len(presentation.slides)
            current_metrics.total_slides = total_slides
            
            if progress_callback and callable(progress_callback):
                progress_callback(0, 0, total_slides)
            
            for slide_idx, slide in enumerate(presentation.slides, 1):
                try:
                    # Update progress
                    if progress_callback and callable(progress_callback):
                        progress = int(slide_idx / total_slides * 100)
                        progress_callback(progress, slide_idx, total_slides)
                    
                    # Process slide
                    slide_data = self._process_slide(slide, filename, slide_idx)
                    
                    if slide_data:
                        dashboard_type = slide_data.get('dashboard_type')
                        if dashboard_type in results:
                            # Validate and clean data
                            validated_data = self.validator.validate_and_clean(
                                slide_data, dashboard_type
                            )
                            if validated_data:
                                results[dashboard_type].append(validated_data)
                                current_metrics.extracted_records += 1
                                current_metrics.confidence_scores.append(
                                    validated_data.get('confidence_score', 0.8)
                                )
                        
                        current_metrics.processed_slides += 1
                    
                except Exception as slide_error:
                    current_metrics.failed_slides += 1
                    self.logger.error(f"Slide {slide_idx} error: {slide_error}")
                    self._log_slide_error(slide_idx, slide_error, ppt_path)
            
            # Calculate data quality score
            if current_metrics.confidence_scores:
                current_metrics.data_quality_score = sum(current_metrics.confidence_scores) / len(current_metrics.confidence_scores)
            
            # Cache results
            try:
                with open(cache_file, 'w', encoding='utf-8') as f:
                    json.dump(results, f, indent=2, default=str)
            except Exception as e:
                self.logger.error(f"Failed to cache results: {e}")
            
            # Update processing time
            current_metrics.processing_time = (datetime.now() - start_time).total_seconds()
            
            # Update instance metrics
            self.metrics = current_metrics
            
            # Log summary
            self._log_processing_summary(ppt_path, results)
            
            return results
            
        except Exception as e:
            self.logger.error(f"Failed to process {ppt_path}: {e}")
            self.logger.error(traceback.format_exc())
            # Return empty results instead of crashing
            return {
                'social_media': [],
                'community_marketing': [],
                'kol_engagement': [],
                'performance_marketing': [],
                'promotion_posts': []
            }
    
    def _process_slide(self, slide, filename: str, slide_idx: int) -> Optional[Dict]:
        """Process individual slide"""
        try:
            slide_text = self._extract_slide_text(slide)
            slide_title = self._get_slide_title(slide)
            
            # Skip if slide is empty
            if not slide_text.strip() and not slide_title.strip():
                return None
            
            # Identify dashboard type with confidence
            dashboard_type, confidence = self._identify_dashboard_type(slide_text, slide_title)
            
            if not dashboard_type:
                return None
            
            # Extract data based on dashboard type
            extraction_methods = {
                'social_media': self._extract_social_media_data,
                'community_marketing': self._extract_community_data,
                'kol_engagement': self._extract_kol_data,
                'performance_marketing': self._extract_performance_data,
                'promotion_posts': self._extract_promotion_data
            }
            
            if dashboard_type in extraction_methods:
                extracted_data = extraction_methods[dashboard_type](
                    slide_text, slide_title, filename, slide_idx
                )
                
                if extracted_data:
                    # Add metadata
                    extracted_data.update({
                        'dashboard_type': dashboard_type,
                        'confidence_score': confidence,
                        'extraction_timestamp': datetime.now().isoformat(),
                        'slide_content_hash': hashlib.md5(slide_text.encode()).hexdigest()[:8]
                    })
                    
                    return extracted_data
            
            return None
            
        except Exception as e:
            self.logger.error(f"Error processing slide {slide_idx}: {e}")
            return None
    
    def _identify_dashboard_type(self, slide_text: str, slide_title: str) -> Tuple[Optional[str], float]:
        """Identify dashboard type with confidence score"""
        slide_content = f"{slide_title}\n{slide_text}".lower()
        
        # More flexible patterns
        dashboard_patterns = {
            'social_media': [
                (r'tf.*value.*mart.*fb.*page', 0.95),
                (r'tf.*value.*mart.*ig.*page', 0.95),
                (r'tf.*value.*mart.*tiktok', 0.95),
                (r'facebook.*page', 0.8),
                (r'instagram.*page', 0.8),
                (r'tiktok.*video', 0.8),
                (r'social.*media', 0.7)
            ],
            'community_marketing': [
                (r'hive.*marketing', 0.95),
                (r'community.*marketing', 0.8),
                (r'group.*performance', 0.7)
            ],
            'kol_engagement': [
                (r'kol.*engagement', 0.95),
                (r'influencer.*collaboration', 0.8),
                (r'creator.*video', 0.7)
            ],
            'performance_marketing': [
                (r'fb.*ads', 0.95),
                (r'ig.*ads', 0.95),
                (r'tiktok.*ads', 0.95),
                (r'ad.*performance', 0.7)
            ],
            'promotion_posts': [
                (r'promotion.*posts', 0.95),
                (r'promotional.*content', 0.8)
            ]
        }
        
        best_match = None
        best_confidence = 0.0
        
        for dashboard_type, patterns in dashboard_patterns.items():
            for pattern, confidence in patterns:
                if re.search(pattern, slide_content, re.IGNORECASE):
                    if confidence > best_confidence:
                        best_match = dashboard_type
                        best_confidence = confidence
                        break
        
        return best_match, best_confidence
    
    def _extract_slide_text(self, slide) -> str:
        """Extract all text from slide including tables"""
        texts = []
        
        # Check if slide has shapes attribute
        if not hasattr(slide, 'shapes'):
            return ""
        
        for shape in slide.shapes:
            try:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    texts.append(clean_text(shape.text))
                
                # Check for tables
                if hasattr(shape, "has_table") and shape.has_table:
                    try:
                        table_text = self._extract_table_text(shape.table)
                        texts.append(table_text)
                    except:
                        pass  # Skip table if error
                        
            except Exception as e:
                self.logger.debug(f"Error processing shape: {e}")
                continue  # Skip this shape and continue
        
        return "\n".join(texts)
    
    def _extract_table_text(self, table) -> str:
        """Extract text from table"""
        rows = []
        for row in table.rows:
            row_text = []
            for cell in row.cells:
                if cell.text.strip():
                    row_text.append(cell.text.strip())
            if row_text:
                rows.append(" | ".join(row_text))
        return "\n".join(rows)
    
    def _get_slide_title(self, slide) -> str:
        """Extract slide title with fallback"""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = clean_text(shape.text)
                if len(text) < 200:  # Likely a title
                    return text
        
        # Check for title placeholder
        for placeholder in slide.placeholders:
            if placeholder.has_text_frame and placeholder.text.strip():
                text = clean_text(placeholder.text)
                if text and len(text) < 200:
                    return text
        
        return f"Slide {slide.slide_id}"
    
    def _extract_social_media_data(self, slide_text: str, slide_title: str,
                                 filename: str, slide_idx: int) -> Dict:
        """Extract social media performance data"""
        platform = self._detect_platform(slide_text, slide_title)
        metrics = self._extract_metrics(slide_text)
        
        data = {
            'platform': platform,
            'post_title': self._extract_title(slide_title),
            'post_date': self._extract_date(slide_text),
            'source_file': filename,
            'slide_number': slide_idx + 1,
            **metrics
        }
        
        # Add platform-specific metadata
        if platform == 'Facebook':
            data['post_type'] = self._detect_post_type(slide_text)
        elif platform == 'Instagram':
            data['content_type'] = self._detect_content_type(slide_text)
        
        return data
    
    def _extract_community_data(self, slide_text: str, slide_title: str,
                              filename: str, slide_idx: int) -> Dict:
        """Extract community marketing data"""
        community_name = self._extract_entity(slide_text, 'community_name')
        
        data = {
            'community_name': community_name,
            'post_title': self._extract_title(slide_title),
            'post_date': self._extract_date(slide_text),
            'source_file': filename,
            'slide_number': slide_idx + 1,
            **self._extract_metrics(slide_text)
        }
        
        return data
    
    def _extract_kol_data(self, slide_text: str, slide_title: str,
                         filename: str, slide_idx: int) -> Dict:
        """Extract KOL engagement data"""
        kol_name = self._extract_entity(slide_text, 'kol_name')
        
        data = {
            'kol_name': kol_name,
            'video_title': self._extract_title(slide_title),
            'video_date': self._extract_date(slide_text),
            'source_file': filename,
            'slide_number': slide_idx + 1,
            **self._extract_metrics(slide_text)
        }
        
        # Classify KOL tier
        views = data.get('views', 0)
        if views >= 1000000:
            data['kol_tier'] = 'Mega'
        elif views >= 100000:
            data['kol_tier'] = 'Macro'
        elif views >= 10000:
            data['kol_tier'] = 'Micro'
        else:
            data['kol_tier'] = 'Nano'
        
        return data
    
    def _extract_performance_data(self, slide_text: str, slide_title: str,
                                 filename: str, slide_idx: int) -> Dict:
        """Extract performance marketing data"""
        platform = self._detect_ad_platform(slide_text)
        ad_type = self._detect_ad_type(slide_text)
        
        data = {
            'platform': platform,
            'ad_type': ad_type,
            'source_file': filename,
            'slide_number': slide_idx + 1,
            **self._extract_metrics(slide_text)
        }
        
        # Calculate derived metrics
        impressions = data.get('impressions', 0)
        if impressions > 0:
            if 'clicks' in data:
                data['ctr'] = (data['clicks'] / impressions) * 100
            
            if 'spend' in data and 'conversions' in data:
                if data['conversions'] > 0:
                    data['cpa'] = data['spend'] / data['conversions']
                    data['roi'] = (data['revenue'] / data['spend']) if data.get('revenue') else None
        
        return data
    
    def _extract_promotion_data(self, slide_text: str, slide_title: str,
                               filename: str, slide_idx: int) -> Dict:
        """Extract promotion posts data"""
        campaign_type = self._detect_campaign_type(slide_text)
        
        data = {
            'campaign_type': campaign_type,
            'post_title': self._extract_title(slide_title),
            'post_date': self._extract_date(slide_text),
            'source_file': filename,
            'slide_number': slide_idx + 1,
            **self._extract_metrics(slide_text)
        }
        
        return data
    
    def _detect_platform(self, text: str, title: str) -> str:
        """Detect social media platform"""
        content = f"{title}\n{text}".lower()
        
        if any(x in content for x in ['fb', 'facebook']):
            return 'Facebook'
        elif any(x in content for x in ['ig', 'instagram']):
            return 'Instagram'
        elif 'tiktok' in content:
            return 'TikTok'
        elif any(x in content for x in ['twitter', 'x.com']):
            return 'Twitter'
        elif any(x in content for x in ['linkedin', 'lnkd']):
            return 'LinkedIn'
        elif 'youtube' in content:
            return 'YouTube'
        
        return 'Unknown'
    
    def _extract_metrics(self, text: str) -> Dict[str, Any]:
        """Extract all metrics from text"""
        metrics = {}
        
        for metric_name, metric_info in self.patterns['metrics'].items():
            match = metric_info['pattern'].search(text)
            if match:
                try:
                    value_str = match.group(1)
                    value = format_number(value_str)
                    if value is not None:
                        metrics[metric_name] = value
                except Exception as e:
                    self.logger.debug(f"Failed to parse {metric_name}: {value_str}, Error: {e}")
        
        return metrics
    
    def _extract_date(self, text: str) -> Optional[str]:
        """Extract date with multiple format support"""
        match = self.patterns['date']['pattern'].search(text)
        if match:
            return match.group(1)
        return None
    
    def _extract_entity(self, text: str, entity_type: str) -> str:
        """Extract named entities"""
        if entity_type in self.patterns['entities']:
            match = self.patterns['entities'][entity_type].search(text)
            if match:
                return clean_text(match.group(1))
        return f"Unknown {entity_type.replace('_', ' ').title()}"
    
    def _extract_title(self, slide_title: str) -> str:
        """Extract and clean title"""
        title = clean_text(slide_title)
        
        # Remove common prefixes
        prefixes = ["Slide", "Performance", "Report", "Dashboard", "Analysis"]
        for prefix in prefixes:
            if title.lower().startswith(prefix.lower()):
                title = title[len(prefix):].strip(" :-\t")
        
        return title[:200] or "Untitled"
    
    def _detect_ad_platform(self, text: str) -> str:
        """Detect advertising platform"""
        text_lower = text.lower()
        
        if any(x in text_lower for x in ['fb ads', 'facebook ads']):
            return 'Facebook Ads'
        elif any(x in text_lower for x in ['ig ads', 'instagram ads']):
            return 'Instagram Ads'
        elif any(x in text_lower for x in ['tiktok ads', 'tiktok advertising']):
            return 'TikTok Ads'
        elif any(x in text_lower for x in ['google ads', 'adwords']):
            return 'Google Ads'
        elif any(x in text_lower for x in ['linkedin ads']):
            return 'LinkedIn Ads'
        
        return 'Unknown Platform'
    
    def _detect_ad_type(self, text: str) -> str:
        """Detect ad type"""
        text_lower = text.lower()
        
        if 'page likes' in text_lower:
            return 'Page Likes'
        elif 'profile visits' in text_lower:
            return 'Profile Visits'
        elif any(x in text_lower for x in ['followers', 'follows']):
            return 'Followers'
        elif 'conversions' in text_lower:
            return 'Conversions'
        elif 'traffic' in text_lower:
            return 'Traffic'
        elif 'awareness' in text_lower:
            return 'Awareness'
        
        return 'Unknown Type'
    
    def _detect_campaign_type(self, text: str) -> str:
        """Detect campaign type"""
        text_lower = text.lower()
        
        if any(x in text_lower for x in ['seasonal', 'holiday']):
            return 'Seasonal'
        elif 'flash sale' in text_lower:
            return 'Flash Sale'
        elif any(x in text_lower for x in ['product launch', 'new product']):
            return 'Product Launch'
        elif 'brand awareness' in text_lower:
            return 'Brand Awareness'
        elif any(x in text_lower for x in ['discount', 'sale', 'promo']):
            return 'Discount/Sale'
        
        return 'General Promotion'
    
    def _detect_post_type(self, text: str) -> str:
        """Detect Facebook post type"""
        text_lower = text.lower()
        
        if 'video' in text_lower:
            return 'Video'
        elif 'image' in text_lower:
            return 'Image'
        elif 'carousel' in text_lower:
            return 'Carousel'
        elif 'story' in text_lower:
            return 'Story'
        elif 'reel' in text_lower:
            return 'Reel'
        
        return 'Post'
    
    def _detect_content_type(self, text: str) -> str:
        """Detect Instagram content type"""
        text_lower = text.lower()
        
        if 'reel' in text_lower:
            return 'Reel'
        elif 'story' in text_lower:
            return 'Story'
        elif 'igtv' in text_lower:
            return 'IGTV'
        elif 'carousel' in text_lower:
            return 'Carousel'
        elif 'feed' in text_lower:
            return 'Feed Post'
        
        return 'Content'
    
    def _log_slide_error(self, slide_idx: int, error: Exception, file_path: str):
        """Log slide processing error"""
        error_log = {
            'timestamp': datetime.now().isoformat(),
            'file': file_path,
            'slide_number': slide_idx,
            'error_type': type(error).__name__,
            'error_message': str(error),
            'stack_trace': traceback.format_exc()
        }
        
        error_file = Path("errors") / f"errors_{datetime.now():%Y%m%d}.json"
        error_file.parent.mkdir(exist_ok=True)
        
        try:
            with open(error_file, 'a', encoding='utf-8') as f:
                json.dump(error_log, f)
                f.write("\n")
        except Exception as e:
            self.logger.error(f"Failed to log error: {e}")
    
    def _log_processing_summary(self, file_path: str, results: Dict):
        """Log processing summary"""
        summary = {
            'timestamp': datetime.now().isoformat(),
            'file': file_path,
            'metrics': {
                'total_slides': self.metrics.total_slides,
                'processed_slides': self.metrics.processed_slides,
                'failed_slides': self.metrics.failed_slides,
                'extracted_records': self.metrics.extracted_records,
                'processing_time_seconds': self.metrics.processing_time,
                'data_quality_score': self.metrics.data_quality_score
            },
            'results_by_type': {
                k: len(v) for k, v in results.items()
            }
        }
        
        self.logger.info(f"Processing summary for {file_path}:")
        self.logger.info(f"  Total slides: {self.metrics.total_slides}")
        self.logger.info(f"  Processed: {self.metrics.processed_slides}")
        self.logger.info(f"  Failed: {self.metrics.failed_slides}")
        self.logger.info(f"  Extracted records: {self.metrics.extracted_records}")
        self.logger.info(f"  Processing time: {self.metrics.processing_time:.2f}s")
        self.logger.info(f"  Data quality: {self.metrics.data_quality_score:.2%}")
        
        for dashboard_type, count in summary['results_by_type'].items():
            if count > 0:
                self.logger.info(f"  {dashboard_type}: {count} records")
    
    def get_processing_metrics(self) -> Dict:
        """Get processing metrics"""
        return {
            'metrics': self.metrics.__dict__,
            'extraction_history': self.extraction_history[-100:]  # Last 100 entries
        }