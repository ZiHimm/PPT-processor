# fixed_simple_processor.py
from pptx import Presentation
import re
import json
import os

class FixedSimpleProcessor:
    """Fixed version with better pattern matching"""
    
    def __init__(self):
        # Use more flexible patterns
        self.patterns = {
            'social_media': [
                # Facebook patterns
                r'tf.*value.*mart.*fb.*page',
                r'tf.*value.*mart.*facebook',
                r'fb.*page.*wallposts',
                r'facebook.*performance',
                
                # Instagram patterns
                r'tf.*value.*mart.*ig.*page',
                r'tf.*value.*mart.*instagram',
                r'ig.*page.*wallposts',
                r'instagram.*performance',
                
                # TikTok patterns
                r'tf.*value.*mart.*tiktok',
                r'tiktok.*video.*performance',
            ],
            
            'community_marketing': [
                r'hive.*marketing',
                r'community.*marketing',
            ],
            
            'kol_engagement': [
                r'kol.*engagement',
                r'influencer.*collaboration',
                r'creator.*video',
            ],
            
            'performance_marketing': [
                r'fb.*ads',
                r'ig.*ads',
                r'tiktok.*ads',
                r'advertising.*performance',
            ],
            
            'promotion_posts': [
                r'promotion.*posts',
                r'promotional.*content',
            ]
        }
        
        # Metrics patterns - more flexible
        self.metrics = {
            'reach': r'reach[:\s]*([\d,\.]+[kmb]?)',
            'views': r'views[:\s]*([\d,\.]+[kmb]?)',
            'engagement': r'engagement[:\s]*([\d,\.]+[kmb]?)',
            'likes': r'likes[:\s]*([\d,\.]+[kmb]?)',
            'shares': r'shares[:\s]*([\d,\.]+[kmb]?)',
            'comments': r'comments[:\s]*([\d,\.]+[kmb]?)',
            'saved': r'saved[:\s]*([\d,\.]+[kmb]?)',
        }
    
    def parse_value(self, text):
        """Convert text like '1.5K' to number"""
        if not text:
            return None
        
        text = str(text).lower().replace(',', '')
        
        if 'k' in text:
            return float(text.replace('k', '')) * 1000
        elif 'm' in text:
            return float(text.replace('m', '')) * 1000000
        elif 'b' in text:
            return float(text.replace('b', '')) * 1000000000
        
        try:
            return float(text)
        except:
            return None
    
    def process_ppt(self, ppt_path, progress_callback=None):
        """Process with better pattern matching"""
        print(f"\nProcessing: {os.path.basename(ppt_path)}")
        print("="*80)
        
        ppt = Presentation(ppt_path)
        total_slides = len(ppt.slides)
        
        results = {
            'social_media': [],
            'community_marketing': [],
            'kol_engagement': [],
            'performance_marketing': [],
            'promotion_posts': []
        }
        
        slide_count = 0
        
        for slide_idx in range(total_slides):
            slide = ppt.slides[slide_idx]
            slide_count += 1
            
            # Update progress
            if progress_callback:
                progress = int((slide_idx + 1) / total_slides * 100)
                progress_callback(progress, slide_idx + 1, total_slides)
            
            # Get all text from slide
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text += shape.text + "\n"
            
            # Convert to lowercase for easier matching
            slide_text_lower = slide_text.lower()
            
            # Check each pattern category
            for category, patterns in self.patterns.items():
                for pattern in patterns:
                    if re.search(pattern, slide_text_lower, re.IGNORECASE):
                        print(f"\n✓ Slide {slide_idx + 1}: Matched {category} pattern: {pattern}")
                        
                        # Extract metrics
                        slide_metrics = {}
                        for metric_name, metric_pattern in self.metrics.items():
                            match = re.search(metric_pattern, slide_text_lower, re.IGNORECASE)
                            if match:
                                value = self.parse_value(match.group(1))
                                if value:
                                    slide_metrics[metric_name] = value
                                    print(f"  Found {metric_name}: {value:,}")
                        
                        # Also look for dates
                        date_match = re.search(r'(\d{1,2}/\d{1,2}/\d{2,4})', slide_text)
                        if date_match:
                            slide_metrics['date'] = date_match.group(1)
                            print(f"  Found date: {slide_metrics['date']}")
                        
                        # Create record if we found anything
                        if slide_metrics:
                            record = {
                                'slide': slide_idx + 1,
                                'category': category,
                                'metrics': slide_metrics,
                                'text_preview': slide_text[:100] + "..."
                            }
                            results[category].append(record)
                            print(f"  ✓ Added to {category}")
                        
                        break  # Stop checking other patterns for this category
        
        # Summary
        print(f"\n{'='*80}")
        print("PROCESSING COMPLETE")
        print("="*80)
        
        total_records = sum(len(r) for r in results.values())
        print(f"\nTotal records extracted: {total_records}")
        
        for category, records in results.items():
            if records:
                print(f"  {category}: {len(records)} records")
        
        return results

def main():
    """Test the fixed processor"""
    ppt_file = "C:/Users/user/Downloads/TFVM_REPORT (FEBRUARY 2025).pptx"
    
    if not os.path.exists(ppt_file):
        print("File not found! Trying to find PPT files...")
        
        # Search for PPT files
        for root, dirs, files in os.walk("."):
            for file in files:
                if file.lower().endswith(('.ppt', '.pptx')):
                    ppt_file = os.path.join(root, file)
                    print(f"Found: {file}")
                    break
            if ppt_file:
                break
    
    print("FIXED SIMPLE PPT PROCESSOR")
    print("="*80)
    print("Using flexible pattern matching (case-insensitive, allows variations)")
    print("="*80)
    
    processor = FixedSimpleProcessor()
    
    # Progress callback
    def show_progress(percent, current, total):
        print(f"Progress: {percent}% ({current}/{total} slides)", end='\r')
    
    # Process
    results = processor.process_ppt(ppt_file, progress_callback=show_progress)
    
    # Save results
    output_file = "fixed_simple_results.json"
    with open(output_file, 'w', encoding='utf-8') as f:
        json.dump(results, f, indent=2, default=str)
    
    print(f"\n✓ Results saved to: {output_file}")
    
    # Show what we found
    if sum(len(r) for r in results.values()) == 0:
        print(f"\n⚠ No data extracted! Let's debug...")
        
        # Show sample of slide content
        print(f"\n{'='*80}")
        print("DEBUG: Checking first 5 slides for content...")
        print("="*80)
        
        ppt = Presentation(ppt_file)
        for i in range(min(5, len(ppt.slides))):
            slide = ppt.slides[i]
            print(f"\nSlide {i+1}:")
            
            # Get all text
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text += shape.text + "\n"
            
            if slide_text.strip():
                print(f"Content (first 200 chars):")
                print(slide_text[:200].replace('\n', ' ') + "...")
                
                # Show lowercase version for pattern matching
                print(f"\nLowercase version:")
                print(slide_text.lower()[:200].replace('\n', ' ') + "...")
            else:
                print("No text found")

if __name__ == "__main__":
    main()
    input("\n\nPress Enter to exit...")