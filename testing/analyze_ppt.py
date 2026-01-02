# ppt_analyzer_complete.py
import os
import sys
from pptx import Presentation
import re
import json
from pathlib import Path

def analyze_ppt_file(ppt_path):
    """Comprehensive analysis of a PPT file"""
    print(f"\n{'='*80}")
    print(f"ANALYZING: {os.path.basename(ppt_path)}")
    print(f"{'='*80}")
    
    try:
        # Load presentation
        presentation = Presentation(ppt_path)
        total_slides = len(presentation.slides)
        print(f"✓ Loaded successfully")
        print(f"Total slides: {total_slides}")
        print(f"File size: {os.path.getsize(ppt_path) / (1024*1024):.1f} MB")
        
        # Dashboard categories and their keywords
        dashboard_categories = {
            'social_media': [
                'facebook', 'fb', 'instagram', 'ig', 'tiktok', 'social media',
                'wallposts', 'post', 'posts', 'reels', 'stories'
            ],
            'community_marketing': [
                'community', 'hive', 'group', 'member', 'membership'
            ],
            'kol_engagement': [
                'kol', 'influencer', 'creator', 'youtuber', 'tiktoker'
            ],
            'performance_marketing': [
                'ads', 'advertising', 'campaign', 'performance', 'roi', 'ctr',
                'cpc', 'conversion', 'impressions', 'clicks'
            ],
            'promotion_posts': [
                'promotion', 'promotional', 'sale', 'discount', 'offer'
            ]
        }
        
        # Metrics to look for
        metrics = [
            'reach', 'views', 'engagement', 'likes', 'shares', 'comments',
            'saved', 'impressions', 'clicks', 'conversions', 'ctr', 'cpc',
            'followers', 'follows', 'page likes', 'profile visits'
        ]
        
        # Initialize counters
        category_counts = {cat: 0 for cat in dashboard_categories}
        metric_counts = {metric: 0 for metric in metrics}
        slide_details = []
        
        print(f"\n{'='*80}")
        print(f"ANALYZING SLIDE CONTENT")
        print(f"{'='*80}")
        
        # Analyze each slide
        for slide_idx in range(total_slides):
            slide = presentation.slides[slide_idx]
            slide_data = {
                'slide_number': slide_idx + 1,
                'text_shapes': 0,
                'tables': 0,
                'charts': 0,
                'categories': [],
                'metrics_found': [],
                'text_preview': ''
            }
            
            # Check for title
            title = ""
            if slide.shapes.title and slide.shapes.title.text:
                title = slide.shapes.title.text.strip()
                slide_data['title'] = title
            
            # Extract all text
            all_text = []
            table_texts = []
            
            for shape in slide.shapes:
                # Text boxes
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text:
                        all_text.append(text)
                        slide_data['text_shapes'] += 1
                
                # Tables
                if hasattr(shape, "has_table") and shape.has_table:
                    slide_data['tables'] += 1
                    try:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                cell_text = cell.text.strip()
                                if cell_text:
                                    row_text.append(cell_text)
                            if row_text:
                                table_texts.append(" | ".join(row_text))
                    except:
                        pass
                
                # Charts (can't extract data, but can note they exist)
                if hasattr(shape, "has_chart") and shape.has_chart:
                    slide_data['charts'] += 1
            
            # Combine all text
            combined_text = "\n".join(all_text + table_texts)
            slide_data['text_preview'] = combined_text[:500] + "..." if len(combined_text) > 500 else combined_text
            
            # Convert to lowercase for searching
            text_lower = combined_text.lower()
            title_lower = title.lower() if title else ""
            
            # Check for dashboard categories
            for category, keywords in dashboard_categories.items():
                for keyword in keywords:
                    if keyword in text_lower or keyword in title_lower:
                        if category not in slide_data['categories']:
                            slide_data['categories'].append(category)
                            category_counts[category] += 1
                        break
            
            # Check for metrics
            for metric in metrics:
                # Look for metric patterns like "Reach: 1,234" or "Engagement 567"
                pattern1 = rf'{metric}[:\s]+([\d,]+(?:\.\d+)?[KMB]?)'
                pattern2 = rf'([\d,]+(?:\.\d+)?[KMB]?)\s+{metric}'
                
                matches1 = re.findall(pattern1, text_lower, re.IGNORECASE)
                matches2 = re.findall(pattern2, text_lower, re.IGNORECASE)
                
                if matches1 or matches2:
                    if metric not in slide_data['metrics_found']:
                        slide_data['metrics_found'].append(metric)
                        metric_counts[metric] += 1
            
            # Only add to details if it has interesting content
            if (slide_data['categories'] or slide_data['metrics_found'] or 
                slide_data['tables'] > 0 or len(combined_text) > 100):
                slide_details.append(slide_data)
        
        # Print summary
        print(f"\n{'='*80}")
        print(f"SUMMARY")
        print(f"{'='*80}")
        
        print(f"\nDashboard Categories Found:")
        for category, count in category_counts.items():
            if count > 0:
                print(f"  ✓ {category}: {count} slides")
            else:
                print(f"  ✗ {category}: Not found")
        
        print(f"\nMetrics Found:")
        for metric, count in metric_counts.items():
            if count > 0:
                print(f"  ✓ {metric}: {count} times")
        
        print(f"\nContent Types:")
        total_text_slides = sum(1 for s in slide_details if s['text_shapes'] > 0)
        total_table_slides = sum(1 for s in slide_details if s['tables'] > 0)
        total_chart_slides = sum(1 for s in slide_details if s['charts'] > 0)
        print(f"  Text-rich slides: {total_text_slides}")
        print(f"  Slides with tables: {total_table_slides}")
        print(f"  Slides with charts: {total_chart_slides}")
        
        # Show detailed slide information
        if slide_details:
            print(f"\n{'='*80}")
            print(f"DETAILED SLIDE ANALYSIS (showing first 10 interesting slides)")
            print(f"{'='*80}")
            
            for i, slide_data in enumerate(slide_details[:10]):
                print(f"\n--- Slide {slide_data['slide_number']} ---")
                if 'title' in slide_data:
                    print(f"Title: {slide_data['title']}")
                
                if slide_data['categories']:
                    print(f"Categories: {', '.join(slide_data['categories'])}")
                
                if slide_data['metrics_found']:
                    print(f"Metrics: {', '.join(slide_data['metrics_found'])}")
                
                if slide_data['tables'] > 0:
                    print(f"Tables: {slide_data['tables']}")
                
                if slide_data['charts'] > 0:
                    print(f"Charts: {slide_data['charts']}")
                
                if slide_data['text_preview']:
                    print(f"Text preview: {slide_data['text_preview'][:200]}...")
        
        # Save analysis to file
        output_file = f"ppt_analysis_{os.path.basename(ppt_path).split('.')[0]}.json"
        analysis_result = {
            'file_info': {
                'filename': os.path.basename(ppt_path),
                'size_mb': os.path.getsize(ppt_path) / (1024*1024),
                'total_slides': total_slides
            },
            'summary': {
                'categories': category_counts,
                'metrics': {k: v for k, v in metric_counts.items() if v > 0},
                'content_types': {
                    'text_slides': total_text_slides,
                    'table_slides': total_table_slides,
                    'chart_slides': total_chart_slides
                }
            },
            'detailed_slides': slide_details[:20]  # First 20 interesting slides
        }
        
        with open(output_file, 'w', encoding='utf-8') as f:
            json.dump(analysis_result, f, indent=2, default=str)
        
        print(f"\n{'='*80}")
        print(f"✓ Analysis saved to: {output_file}")
        print(f"{'='*80}")
        
        return analysis_result
        
    except Exception as e:
        print(f"\n✗ Error analyzing file: {e}")
        import traceback
        traceback.print_exc()
        return None

def check_processor_config():
    """Check if processor configuration matches PPT content"""
    config_file = "config/dashboard_config.yaml"
    
    if os.path.exists(config_file):
        print(f"\n{'='*80}")
        print(f"CHECKING PROCESSOR CONFIGURATION")
        print(f"{'='*80}")
        
        import yaml
        with open(config_file, 'r', encoding='utf-8') as f:
            config = yaml.safe_load(f)
        
        if 'dashboard_config' in config:
            config = config['dashboard_config']
            
            print(f"\nCurrent configuration:")
            for category in ['social_media', 'community_marketing', 'kol_engagement', 
                            'performance_marketing', 'promotion_posts']:
                if category in config:
                    print(f"\n{category.upper()}:")
                    if 'keywords' in config[category]:
                        print(f"  Keywords: {config[category]['keywords']}")
                    if 'metrics' in config[category]:
                        print(f"  Metrics: {config[category]['metrics']}")
    else:
        print(f"\n⚠ Configuration file not found: {config_file}")

def main():
    """Main function"""
    print("PPT CONTENT ANALYZER")
    print("="*80)
    
    # Find PPT files
    ppt_files = []
    current_dir = os.getcwd()
    
    # Look for PPT files in common locations
    search_paths = [
        current_dir,
        os.path.join(current_dir, "downloads"),
        os.path.join(current_dir, "data"),
        os.path.join(os.path.expanduser("~"), "Downloads"),
    ]
    
    for path in search_paths:
        if os.path.exists(path):
            for ext in ['.pptx', '.ppt']:
                for file in Path(path).glob(f"*{ext}"):
                    ppt_files.append(str(file))
    
    if not ppt_files:
        print("No PPT files found!")
        print("\nPlease specify a PPT file path:")
        ppt_path = input("Path to PPT file: ").strip().strip('"')
        if os.path.exists(ppt_path):
            ppt_files.append(ppt_path)
        else:
            print(f"File not found: {ppt_path}")
            return
    
    print(f"\nFound {len(ppt_files)} PPT file(s):")
    for i, file in enumerate(ppt_files, 1):
        print(f"{i}. {os.path.basename(file)}")
    
    # Analyze each file
    for i, ppt_file in enumerate(ppt_files, 1):
        print(f"\n{'='*80}")
        print(f"Analyzing file {i}/{len(ppt_files)}")
        analysis = analyze_ppt_file(ppt_file)
        
        if analysis and i < len(ppt_files):
            input("\nPress Enter to continue to next file...")
    
    # Check processor configuration
    check_processor_config()
    
    print(f"\n{'='*80}")
    print("ANALYSIS COMPLETE")
    print("="*80)
    print("\nRecommendations based on analysis:")
    print("1. If slides contain TABLES: The processor needs table extraction")
    print("2. If slides contain CHARTS: Data cannot be automatically extracted")
    print("3. If keywords not found: Update config/dashboard_config.yaml")
    print("4. If metrics in wrong format: Update regex patterns in processor")
    
    input("\nPress Enter to exit...")

if __name__ == "__main__":
    try:
        main()
    except KeyboardInterrupt:
        print("\n\nAnalysis interrupted by user.")
    except Exception as e:
        print(f"\n\nUnexpected error: {e}")
        import traceback
        traceback.print_exc()
        input("\nPress Enter to exit...")