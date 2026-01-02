# test_ppt_speed.py
import os
import time
from pptx import Presentation

ppt_file = "C:\\Users\\user\\Downloads\\TFVM_REPORT (FEBRUARY 2025).pptx"

print(f"Testing PPT file: {ppt_file}")
print(f"File size: {os.path.getsize(ppt_file) / (1024*1024):.1f} MB")

try:
    start = time.time()
    print("Loading presentation...")
    presentation = Presentation(ppt_file)
    load_time = time.time() - start
    print(f"✓ Loaded in {load_time:.1f} seconds")
    
    # Correct way to get number of slides
    num_slides = len(presentation.slides)
    print(f"Number of slides: {num_slides}")
    
    # Test processing first few slides - CORRECT WAY
    print("\nProcessing first 5 slides...")
    slide_count = 0
    
    # Method 1: Using enumerate directly on presentation.slides
    for i, slide in enumerate(presentation.slides):
        if slide_count >= 5:
            break
            
        slide_start = time.time()
        try:
            shapes = len(slide.shapes)
            # Try to extract some text
            text_content = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_content += shape.text[:50] + "..."
            
            slide_time = time.time() - slide_start
            print(f"  Slide {i+1}: {shapes} shapes, {len(text_content)} chars, processed in {slide_time:.3f}s")
            
        except Exception as e:
            print(f"  Slide {i+1}: Error - {str(e)[:50]}")
        
        slide_count += 1
    
    total_time = time.time() - start
    print(f"\nTotal test time: {total_time:.1f} seconds")
    
except Exception as e:
    print(f"✗ Error: {e}")
    import traceback
    traceback.print_exc()